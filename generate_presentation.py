from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import pandas as pd

# Couleurs
VERT_FONCE = RGBColor(0x2C, 0x5F, 0x2D)
VERT_CLAIR = RGBColor(0x97, 0xBC, 0x62)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
GRIS = RGBColor(0x44, 0x44, 0x44)
ROUGE = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_slide(prs):
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)

def add_rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=BLANC, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

# ══ SLIDE 1 — COVER ════════════════════════════════════════════════
s1 = add_slide(prs)
add_rect(s1, 0, 0, 13.33, 7.5, VERT_FONCE)
add_rect(s1, 0, 0, 0.2, 7.5, VERT_CLAIR)
add_rect(s1, 0.2, 5.5, 13.13, 0.08, VERT_CLAIR)
add_text(s1, "GreenFinance Europe", 0.5, 1.2, 12, 1.5, size=44, bold=True)
add_text(s1, "Analyse de la transition énergétique européenne 2007-2021", 0.5, 2.5, 12, 0.8, size=20)
add_text(s1, "Source : Banque mondiale  ·  Python · Streamlit · Pandas", 0.5, 3.2, 12, 0.6, size=14, color=VERT_CLAIR)
add_rect(s1, 0.5, 4.0, 2.8, 1.0, RGBColor(0x1A, 0x3A, 0x1A))
add_rect(s1, 3.5, 4.0, 2.8, 1.0, RGBColor(0x1A, 0x3A, 0x1A))
add_rect(s1, 6.5, 4.0, 2.8, 1.0, RGBColor(0x1A, 0x3A, 0x1A))
add_rect(s1, 9.5, 4.0, 2.8, 1.0, RGBColor(0x1A, 0x3A, 0x1A))
add_text(s1, "10", 0.5, 4.05, 2.8, 0.5, size=28, bold=True, color=VERT_CLAIR, align=PP_ALIGN.CENTER)
add_text(s1, "Pays analysés", 0.5, 4.55, 2.8, 0.35, size=11, color=BLANC, align=PP_ALIGN.CENTER)
add_text(s1, "15", 3.5, 4.05, 2.8, 0.5, size=28, bold=True, color=VERT_CLAIR, align=PP_ALIGN.CENTER)
add_text(s1, "Années de données", 3.5, 4.55, 2.8, 0.35, size=11, color=BLANC, align=PP_ALIGN.CENTER)
add_text(s1, "42.5%", 6.5, 4.05, 2.8, 0.5, size=28, bold=True, color=VERT_CLAIR, align=PP_ALIGN.CENTER)
add_text(s1, "Objectif UE 2030", 6.5, 4.55, 2.8, 0.35, size=11, color=BLANC, align=PP_ALIGN.CENTER)
add_text(s1, "0.85", 9.5, 4.05, 2.8, 0.5, size=28, bold=True, color=VERT_CLAIR, align=PP_ALIGN.CENTER)
add_text(s1, "Corrélation", 9.5, 4.55, 2.8, 0.35, size=11, color=BLANC, align=PP_ALIGN.CENTER)
add_text(s1, "Rafika Cervera  ·  github.com/RAFCERAY  ·  Mars 2026", 0.5, 6.9, 12, 0.4, size=10, color=RGBColor(0x88,0x88,0x88))

# ══ SLIDE 2 — CONTEXTE ══════════════════════════════════════════════
s2 = add_slide(prs)
add_rect(s2, 0, 0, 13.33, 0.9, VERT_FONCE)
add_text(s2, "Contexte & Données", 0.4, 0.1, 12, 0.7, size=26, bold=True)
add_rect(s2, 0.4, 1.0, 5.8, 5.8, RGBColor(0xF5,0xF5,0xF5))
add_rect(s2, 0.4, 1.0, 5.8, 0.08, VERT_CLAIR)
add_text(s2, "Mission", 0.5, 1.1, 5.6, 0.4, size=14, bold=True, color=VERT_FONCE)
add_text(s2, "Analyser la progression des énergies renouvelables dans 10 pays européens et mesurer l'écart avec l'objectif UE 2030 de 42.5%.", 0.5, 1.55, 5.6, 1.2, size=12, color=GRIS)
add_text(s2, "Datasets", 0.5, 2.8, 5.6, 0.4, size=14, bold=True, color=VERT_FONCE)
add_text(s2, "renouvelables_clean.csv\n150 lignes · 2007-2021\nIndicateur EG.FEC.RNEW.ZS", 0.5, 3.25, 5.6, 0.9, size=11, color=GRIS)
add_text(s2, "electricite_clean.csv\n120 lignes · 2010-2021\nIndicateur EG.ELC.RNEW.ZS", 0.5, 4.2, 5.6, 0.9, size=11, color=GRIS)
add_text(s2, "Qualité : 0 doublons · 0 valeurs manquantes", 0.5, 5.2, 5.6, 0.4, size=11, bold=True, color=VERT_FONCE)
add_rect(s2, 6.5, 1.0, 6.4, 5.8, RGBColor(0xF5,0xF5,0xF5))
add_rect(s2, 6.5, 1.0, 6.4, 0.08, VERT_CLAIR)
add_text(s2, "Pays analysés", 6.6, 1.1, 6.2, 0.4, size=14, bold=True, color=VERT_FONCE)
pays = ["Austria · Central", "Belgium · Western", "Denmark · Nordic",
        "France · Western", "Germany · Central", "Italy · Southern",
        "Netherlands · Western", "Poland · Eastern", "Spain · Southern", "Sweden · Nordic"]
for i, p in enumerate(pays):
    y = 1.6 + i * 0.5
    add_rect(s2, 6.6, y, 5.8, 0.38, VERT_FONCE if "Nordic" in p else RGBColor(0xE8,0xF5,0xE9))
    color = BLANC if "Nordic" in p else VERT_FONCE
    add_text(s2, p, 6.7, y+0.04, 5.6, 0.3, size=11, color=color)

# ══ SLIDE 3 — RÉSULTATS CLÉS ════════════════════════════════════════
s3 = add_slide(prs)
add_rect(s3, 0, 0, 13.33, 0.9, VERT_FONCE)
add_text(s3, "Résultats clés", 0.4, 0.1, 12, 0.7, size=26, bold=True)
kpis = [
    ("14.1% → 24.3%", "Progression EU\n2007-2021", VERT_FONCE),
    ("+10.2 pts", "En 14 ans", VERT_CLAIR),
    ("1 pays / 10", "Objectif atteint\n(Sweden)", ORANGE),
    ("+18.2 pts", "Écart restant\navant 2030", ROUGE),
]
for i, (val, label, color) in enumerate(kpis):
    x = 0.4 + i * 3.2
    add_rect(s3, x, 1.1, 3.0, 2.0, color)
    add_text(s3, val, x, 1.3, 3.0, 0.9, size=22, bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, label, x, 2.2, 3.0, 0.7, size=12, align=PP_ALIGN.CENTER)
add_text(s3, "Classement 2021", 0.4, 3.3, 12, 0.4, size=16, bold=True, color=VERT_FONCE)
pays_data = [
    ("🥇 Sweden", "57.9%", "✅ Objectif atteint", VERT_FONCE),
    ("🥈 Denmark", "39.5%", "🟡 Proche", VERT_CLAIR),
    ("🥉 Austria", "36.0%", "🟡 Proche", VERT_CLAIR),
    ("France", "16.2%", "🔴 En retard", ROUGE),
    ("Belgium", "11.7%", "🔴 Très en retard", ROUGE),
]
for i, (pays, val, statut, color) in enumerate(pays_data):
    x = 0.4 + i * 2.55
    add_rect(s3, x, 3.8, 2.4, 1.8, color)
    add_text(s3, pays, x, 3.85, 2.4, 0.5, size=12, bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, val, x, 4.3, 2.4, 0.5, size=20, bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, statut, x, 4.85, 2.4, 0.5, size=10, align=PP_ALIGN.CENTER)

# ══ SLIDE 4 — ANALYSE RÉGIONALE ═════════════════════════════════════
s4 = add_slide(prs)
add_rect(s4, 0, 0, 13.33, 0.9, VERT_FONCE)
add_text(s4, "Analyse par région", 0.4, 0.1, 12, 0.7, size=26, bold=True)
regions = [
    ("Nordic", "~39%", "Denmark + Sweden\nModèle européen", VERT_FONCE),
    ("Central", "~23%", "Austria + Germany\nProgression correcte", VERT_CLAIR),
    ("Southern", "~15%", "Spain + Italy\nPotentiel solaire élevé", ORANGE),
    ("Eastern", "~11%", "Poland\nEn retard", RGBColor(0xE6,0x7E,0x22)),
    ("Western", "~9%", "France · Belgium · NL\nLe plus en retard", ROUGE),
]
for i, (region, val, desc, color) in enumerate(regions):
    x = 0.3 + i * 2.6
    add_rect(s4, x, 1.1, 2.4, 4.5, color)
    add_text(s4, region, x, 1.2, 2.4, 0.5, size=14, bold=True, align=PP_ALIGN.CENTER)
    add_text(s4, val, x, 1.8, 2.4, 0.8, size=28, bold=True, align=PP_ALIGN.CENTER)
    add_text(s4, desc, x, 2.7, 2.4, 1.0, size=11, align=PP_ALIGN.CENTER)
    add_rect(s4, x, 3.8, 2.4, 0.08, BLANC)
add_text(s4, "Objectif UE 2030 : 42.5% — Aucune région n'a encore atteint cet objectif",
         0.4, 5.8, 12.5, 0.5, size=13, bold=True, color=ROUGE, align=PP_ALIGN.CENTER)

# ══ SLIDE 5 — DASHBOARD ══════════════════════════════════════════════
s5 = add_slide(prs)
add_rect(s5, 0, 0, 13.33, 0.9, VERT_FONCE)
add_text(s5, "Dashboard Streamlit", 0.4, 0.1, 12, 0.7, size=26, bold=True)
pages = [
    ("🌿", "Énergies Renouvelables", "Évolution par pays · Comparaison · Données brutes"),
    ("💚", "Électricité Renouvelable", "Part électricité verte · Tendances · Par région"),
]
for i, (icon, title, desc) in enumerate(pages):
    y = 1.2 + i * 2.5
    add_rect(s5, 0.4, y, 12.5, 2.2, RGBColor(0xF5,0xF5,0xF5))
    add_rect(s5, 0.4, y, 12.5, 0.08, VERT_CLAIR)
    add_text(s5, icon + "  " + title, 0.6, y+0.15, 8, 0.5, size=18, bold=True, color=VERT_FONCE)
    add_text(s5, desc, 0.6, y+0.7, 8, 0.5, size=13, color=GRIS)
add_text(s5, "Stack technique", 0.4, 6.2, 12, 0.4, size=14, bold=True, color=VERT_FONCE)
add_text(s5, "Python  ·  Streamlit  ·  Plotly  ·  Pandas  ·  World Bank API  ·  GitHub",
         0.4, 6.6, 12, 0.4, size=13, color=GRIS)

# ══ SLIDE 6 — RECOMMANDATIONS ════════════════════════════════════════
s6 = add_slide(prs)
add_rect(s6, 0, 0, 13.33, 0.9, VERT_FONCE)
add_text(s6, "Recommandations & Conclusion", 0.4, 0.1, 12, 0.7, size=26, bold=True)
recs = [
    ("🏆", "S'inspirer du modèle danois", "Denmark a gagné +21.8 pts en 14 ans grâce à une politique énergétique volontariste. Investissement massif dans l'éolien offshore."),
    ("⚡", "Accélérer en Europe occidentale", "France, Belgium et Netherlands doivent tripler leur rythme de transition pour espérer atteindre l'objectif 2030."),
    ("☀️", "Exploiter le potentiel solaire du Sud", "Spain et Italy ont un potentiel solaire énorme encore sous-exploité. Des investissements ciblés pourraient accélérer leur transition."),
]
for i, (icon, title, desc) in enumerate(recs):
    y = 1.1 + i * 1.8
    add_rect(s6, 0.4, y, 12.5, 1.6, RGBColor(0xF5,0xF5,0xF5))
    add_rect(s6, 0.4, y, 0.08, 1.6, VERT_CLAIR)
    add_text(s6, icon + "  " + title, 0.6, y+0.1, 11, 0.5, size=15, bold=True, color=VERT_FONCE)
    add_text(s6, desc, 0.6, y+0.65, 11.5, 0.8, size=12, color=GRIS)
add_rect(s6, 0.4, 6.5, 12.5, 0.7, VERT_FONCE)
add_text(s6, "Seule la Suède a atteint l'objectif. Le Danemark est le modèle. L'Europe doit accélérer.",
         0.4, 6.55, 12.5, 0.5, size=14, bold=True, align=PP_ALIGN.CENTER)

prs.save("reports/GreenFinance_Presentation.pptx")
print("✅ Présentation créée : reports/GreenFinance_Presentation.pptx")